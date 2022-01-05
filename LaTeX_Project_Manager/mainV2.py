# -*- coding: utf-8 -*-
##########################################################
# PACKAGES
##########################################################
import os # Pour la manipulation des dossiers
import shutil # Pour la manipulation des fichiers
from datetime import datetime # Pour la création de l'historique
import traceback # Pour la gestion des erreurs
from pptx import Presentation # Pour la création du .pptx
##########################################################
# CODE
##########################################################
# Définitions des variables globales :
now = datetime.now() # Récupération de la date et de l'heure
date_time = now.strftime("%m/%d/%Y, %H:%M:%S") # Mise en forme
# Dossier où sont stockés les projets et documents à créer --> USER INPUT :
main_folder_path = "C:\\Users\\Lorinfo\\Documents\\GitHub\\SUPAERO"
# Dossier où sont stockés les fichiers sources TeX --> USER INPUT :
source_path = "C:\\Users\\Lorinfo\\Documents\\GitHub\\LaTeX-Project-Manager\\LaTeX_Project_Manager\\source_TeX"
##########################################################
# Définitions des fonctions utiles :
# --------------------------------------------------------
# Creation d'un historique
def historique(folder):
    """
    Initialisation de l'historique
    """
    fname = folder + "\\historique.txt"
    with open(fname, "w") as file:
        file.write("CREATION HISTORIQUE - INITIALISATION A {date}\n".format(date=date_time))
        file.write("===========================================================\n")
# Ajout d'une entrée dans l'historique
def new_entry(project_name,folder):
    """
    Ajout d'un élément à l'historique
    """
    fname = folder + "\\historique.txt"
    with open(fname,"a") as file:
        file.write("Ajout d'un élément à l'historique à {date}\n".format(date=date_time))
        file.write("--> CREATION DE {pname}\n".format(pname=project_name))
        file.write("-------------------------------------------\n")

# Création du répertoire   
def make_folder(cd,fdname):
    """
    Fonction permettant de créer des dossiers dans un répertoire donné.
    cd --> str ; working directory utilisé pour stocker les fichiers
    fdname --> str ; nom du dossier à créer
    """
    folder_name = fdname
    folder_path = cd + "\\{name}".format(name=folder_name)
    try:
        os.makedirs(folder_path)    
        print("Directory ",folder_name," Created ")
    except FileExistsError:
        print("Directory ",folder_name," already exists")

# Récupération des fichiers sources 
def copytex(fname,fdname):
    """
    Fonction permettant de copier des fichiers d'un répertoire à un autre
    fname --> str ; nom du fichier à copier
    fdname --> str ; nom du dossier dans lequel envoyer le fichier
    NOTE : LE FICHIER PAR DEFAUT DANS LEQUEL SONT PRIS LES FICHIERS .TEX EST source_TeX
    """
    original = source_path + "\\{name}".format(name=fname)
    target = main_folder_path + "\\{fname}".format(fname=fdname) + "\\{name}".format(name=fname)
    shutil.copyfile(original, target)

# Création du .pptx
def make_pptx(project_name):
    """
    Fonction permettant de créer la présentation PowerPoint pour stockr
    les images du document
    project_name --> str ; nom du projet où on souhaite créer le fichier
    """
    print("Fonction executée")
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = project_name
    subtitle.text = "Illustrations"
    prs.save(main_folder_path + "\\{0}".format(project_name) +"\\images" + "\\images.pptx")
    # On sauvegarde le document dans le folder 'images'.

# Routine de creation
def create_project(project_name,classe,images=True,pptx=False):
    """
    Routine de création de projet LaTeX.
    project_name --> str ; nom du projet à créer
    classe --> str ; classe du document à créer
    images --> bool, default = True ; indique s'il faut créer un dossier pour stocker des images
    pptx --> bool, default = False ; indique s'il faut créer dans 'images' un .pptx
    """
    marker = True
    cd = main_folder_path # working directory où il faut importer les dossiers
    path = cd + "\\{fdname}".format(fdname=project_name) # Path du nouveau projet à setup
    new_entry(project_name,main_folder_path) # Ajout d'un élément à l'historique correspondant au nouveau projet
    list_file = os.listdir(source_path) # Liste des fichiers se trouvant dans le dossier source
    if classe == "fiche":
        make_folder(cd,project_name)
        for file in list_file:
            if file.startswith("bristol") or file.startswith("mypackage"):
                copytex(file,project_name)
        if images:
            make_folder(path,"images")
            if pptx:
                make_pptx(project_name)
    elif classe == "standard": # C'est-à-dire 'article'
        make_folder(cd,project_name)
        for file in list_file:
            if file.startswith("glob") or file.startswith("standard"):
                copytex(file,project_name)
        if images: # S'il faut créer un dossier contenant les images du projet
            make_folder(path,"images")
            if pptx:
                make_pptx(project_name)
    elif classe == "book": # C'est-à-dire un projet plus complet
        make_folder(cd,project_name)
        for file in list_file:
            if file.startswith("glob") or file.startswith("standard"):
                    copytex(file,project_name)
            if images:
                make_folder(path,"images")
                if pptx:
                    make_pptx(project_name)
    elif classe == "beam":
        make_folder(cd,project_name)
        for file in list_file:
            if file.startswith("glob") or file.startswith("beam"):
                copytex(file,project_name)
        if images:
            make_folder(path,"images")
            # Pas besoin de créer un fichier .pptx ici
    else:
        marker = False
        print("Aucune correspondance trouvée. Aucun dossier généré.")
    if marker:
        try:
            os.chdir(path) # Changement de working directory
            # print("Current working directory: {0}".format(os.getcwd())) # Debug
        except FileNotFoundError:
            print("Le dossier {0} n'existe pas.".format(path))
        except NotADirectoryError:
            print("{0} n'est pas un dossier.".format(path))
        except PermissionError:
            print("Vous n'avez pas la permission de changer {0}".format(path))
        list_file = os.listdir(os.getcwd())
        # print(list_file) # Debug
        if classe in ["book","standard"]: # Pour fiche, pas besoin
            # J'utilise ici mes conventions : elles sont prises en compte dans les fichiers .tex
            os.rename("glob_macros.tex","macros.tex")
            os.rename("glob_PackagesGlobStandard.tex","packages.tex")
            os.rename("standard_snippet_template.tex","main.tex")
        if classe == "beam":
            os.rename("glob_macros.tex","macros.tex")
            os.rename("glob_PackagesGlobStandard.tex","packages.tex")
            os.rename("beam_snippet_template_beamer.tex","main.tex")
        os.startfile(path) # Pour ouvrir quand il a été créé.
    else:
        print("Aucun dossier n'a été créé. Fin du programme.")

def remove_project(project_name):
    """
    Routine de suppression d'un projet
    project_name --> str ; nom du projet à supprimer.
    """
    delpath = main_folder_path + "\\{0}".format(project_name)
    try:
        shutil.rmtree(delpath)
    except:
        traceback.print_exc() # Pour récuperer les erreurs et les informations correspondantes

##################################################################
# INITIALISATION D'UN NOUVEAU PROJET
# historique(main_folder_path)

create_project("PC_MecaFluChoc_1","standard",images=True,pptx=True)
#remove_project("test")